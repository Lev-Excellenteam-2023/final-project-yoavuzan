import argparse
import asyncio
import concurrent.futures
import json
import openai
from pptx import Presentation

openai.api_key = "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" # Replace with your API key


def extract_text_from_slide(slide):
    # Extract text from a slide
    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text + " "
    return text.strip()


async def complete_chat(presentation_text):
    # Complete a chat with GPT-3
    executor = concurrent.futures.ThreadPoolExecutor()
    # Set the executor as the default executor for asyncio
    loop = asyncio.get_event_loop()
    loop.set_default_executor(executor)
    text = presentation_text + "can you write me a short paragraph about these text, what are they talking about?"
    response = await loop.run_in_executor(None, lambda: openai.ChatCompletion.create
    (model="gpt-3.5-turbo",
     messages=[{"role": "user", "content": text}],
     timeout=100000 # seconds
     )) # Run the API request in a separate thread
    assistant_reply = response.choices[0].message
    assistant_reply_text = assistant_reply.content
    return assistant_reply_text


async def responses_from_server(presentation_path):
    # Extract text from a presentation and complete a chat with GPT-3 for each slide
    # get path to presentation file from command line arguments  and open it
    presentation = Presentation(presentation_path)
    presentation_name = presentation_path.split("/")[-1].split(".")[0]
    tasks = []
    # Create a task for each slide
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = extract_text_from_slide(slide)
        if slide_text: # If the slide has text
            task = asyncio.create_task(complete_chat(slide_text)) # Create a task for the slide
            tasks.append((index, slide_text, task))

    responses = []
    response_dict = {}
    # Wait for the tasks to complete
    for index, slide_text, task in tasks:
        try:
            response = await task
            responses.append(response)
        except Exception as e:
            response = f"Error occurred while processing slide {index}. Error message: {str(e)}"
        response_dict[f"response {index}"] = {"text": slide_text, "response": response}

    return response_dict


async def main(presentation_path):
    # Log the start of the main function
    print("Main function started.")

    responses = await responses_from_server(presentation_path)

    # Log the responses
    for index, response in enumerate(responses):
        print(f"Response {index}: {response}")

    # Save responses to a JSON file
    with open("responses.json", "w") as f:
        json.dump(responses, f)

    print("The JSON file with responses has been saved to the current directory.")

    # Log the end of the main function
    print("Main function finished.")


if __name__ == "__main__":
    # Create an argument parser
    parser = argparse.ArgumentParser(description="Extract text from a PowerPoint presentation and generate responses using GPT-3.5 Turbo.")
    # Add an argument for the presentation file path
    parser.add_argument("presentation", type=str, help="Path to the PowerPoint presentation file")
    # Parse the arguments
    args = parser.parse_args()

    # Create a ThreadPoolExecutor
    executor = concurrent.futures.ThreadPoolExecutor()

    # Set the executor as the default executor for asyncio
    loop = asyncio.get_event_loop()
    loop.set_default_executor(executor)

    # Log the start of the program
    print("Program started.")

    # Run the main async function
    loop.run_until_complete(main(args.presentation))

    # Shutdown the executor
    executor.shutdown()

    # Log the end of the program
    print("Program finished.")
